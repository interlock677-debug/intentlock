#!/usr/bin/env python3
"""
IntentLock v4.0 Pitch Deck Generator
Creates an 8-slide presentation with modern enterprise design
"""

import io

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Color Palette
DARK_BG = RGBColor(15, 18, 25)  # Deep dark blue-black
NEON_CYAN = RGBColor(0, 229, 255)  # Bright cyan
NEON_AMBER = RGBColor(255, 191, 0)  # Bright amber
TEXT_PRIMARY = RGBColor(240, 241, 243)  # Off-white
TEXT_SECONDARY = RGBColor(180, 190, 200)  # Light gray
ACCENT_DARK = RGBColor(30, 40, 55)  # Dark blue accent

def create_blank_image(width, height, color):
    """Create a blank colored image"""
    return Image.new('RGB', (width, height), color)

def create_wireframe_background():
    """Create abstract architectural wireframe background"""
    width, height = 1920, 1440
    img = Image.new('RGB', (width, height), DARK_BG)
    draw = ImageDraw.Draw(img, 'RGBA')

    # Draw grid lines with low opacity
    line_color = (0, 229, 255, 20)  # Cyan with transparency

    for i in range(0, width, 120):
        draw.line([(i, 0), (i, height)], fill=line_color, width=1)

    for i in range(0, height, 120):
        draw.line([(0, i), (width, i)], fill=line_color, width=1)

    # Draw some geometric shapes for visual interest
    shape_color = (255, 191, 0, 15)  # Amber with transparency
    draw.rectangle([(100, 100), (400, 400)], outline=shape_color, width=2)
    draw.rectangle([(1400, 900), (1800, 1200)], outline=shape_color, width=2)

    # Save to bytes
    img_byte_arr = io.BytesIO()
    img.save(img_byte_arr, format='PNG')
    img_byte_arr.seek(0)
    return img_byte_arr

def create_comparison_diagram():
    """Create side-by-side comparison diagram"""
    width, height = 1920, 1080
    img = Image.new('RGB', (width, height), DARK_BG)
    draw = ImageDraw.Draw(img, 'RGBA')

    # Left side - Fixed TTL (red/bad)
    draw.rectangle([(50, 50), (900, 1000)], outline=(255, 100, 100, 200), width=3)

    # Right side - Adaptive Nonce (cyan/good)
    draw.rectangle([(1000, 50), (1850, 1000)], outline=(0, 229, 255, 200), width=3)

    img_byte_arr = io.BytesIO()
    img.save(img_byte_arr, format='PNG')
    img_byte_arr.seek(0)
    return img_byte_arr

def create_cache_flowchart():
    """Create L1/L2 cache flowchart"""
    width, height = 1920, 1080
    img = Image.new('RGB', (width, height), DARK_BG)
    draw = ImageDraw.Draw(img, 'RGBA')

    # L1 Local Cache box
    draw.rectangle([(200, 200), (800, 500)], fill=(0, 229, 255, 30), outline=(0, 229, 255, 200), width=3)

    # L2 Redis box
    draw.rectangle([(200, 600), (800, 900)], fill=(255, 191, 0, 30), outline=(255, 191, 0, 200), width=3)

    # Arrow between them
    draw.line([(500, 500), (500, 600)], fill=(0, 229, 255, 150), width=4)

    img_byte_arr = io.BytesIO()
    img.save(img_byte_arr, format='PNG')
    img_byte_arr.seek(0)
    return img_byte_arr

def create_timeline_diagram():
    """Create blocking vs async timeline"""
    width, height = 1920, 1080
    img = Image.new('RGB', (width, height), DARK_BG)
    draw = ImageDraw.Draw(img, 'RGBA')

    # Blocking timeline (top)
    draw.rectangle([(200, 150), (1700, 350)], outline=(255, 100, 100, 200), width=2)
    draw.rectangle([(300, 200), (1200, 300)], fill=(255, 100, 100, 80))

    # Async timeline (bottom)
    draw.rectangle([(200, 650), (1700, 850)], outline=(0, 229, 255, 200), width=2)
    draw.rectangle([(300, 700), (900, 800)], fill=(0, 229, 255, 80))
    draw.rectangle([(1000, 700), (1600, 800)], fill=(255, 191, 0, 80))

    img_byte_arr = io.BytesIO()
    img.save(img_byte_arr, format='PNG')
    img_byte_arr.seek(0)
    return img_byte_arr

def create_velocity_graph():
    """Create velocity detection graph"""
    width, height = 1920, 1080
    img = Image.new('RGB', (width, height), DARK_BG)
    draw = ImageDraw.Draw(img, 'RGBA')

    # Axes
    draw.line([(200, 900), (1700, 900)], fill=(0, 229, 255, 150), width=3)  # X-axis
    draw.line([(200, 100), (200, 900)], fill=(0, 229, 255, 150), width=3)   # Y-axis

    # Single request curve (flat)
    points_single = [(300, 800), (700, 750), (1100, 700), (1500, 650)]
    for i in range(len(points_single)-1):
        draw.line([points_single[i], points_single[i+1]], fill=(255, 100, 100, 150), width=3)

    # Cumulative curve (steep)
    points_cumulative = [(300, 850), (700, 600), (1100, 300), (1500, 150)]
    for i in range(len(points_cumulative)-1):
        draw.line([points_cumulative[i], points_cumulative[i+1]], fill=(0, 229, 255, 200), width=4)

    img_byte_arr = io.BytesIO()
    img.save(img_byte_arr, format='PNG')
    img_byte_arr.seek(0)
    return img_byte_arr

def create_kms_diagram():
    """Create HSM/KMS security diagram"""
    width, height = 1920, 1080
    img = Image.new('RGB', (width, height), DARK_BG)
    draw = ImageDraw.Draw(img, 'RGBA')

    # Server process box
    draw.rectangle([(300, 200), (900, 500)], outline=(255, 100, 100, 200), width=3)

    # HSM/KMS box
    draw.rectangle([(1100, 200), (1700, 500)], outline=(0, 229, 255, 200), width=3)

    # Arrow with "Envelope Signing"
    draw.line([(900, 350), (1100, 350)], fill=(0, 229, 255, 150), width=4)

    img_byte_arr = io.BytesIO()
    img.save(img_byte_arr, format='PNG')
    img_byte_arr.seek(0)
    return img_byte_arr

def set_background(slide, color):
    """Set solid background color for a slide"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_slide(prs):
    """Slide 1: Title Slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    set_background(slide, DARK_BG)

    # Add wireframe background image
    img = create_wireframe_background()
    left = top = Inches(0)
    pic = slide.shapes.add_picture(img, left, top, width=prs.slide_width, height=prs.slide_height)

    # Move picture to back
    slide.shapes._spTree.remove(pic._element)
    slide.shapes._spTree.insert(2, pic._element)

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(2))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True

    p = title_frame.paragraphs[0]
    p.text = "IntentLock v4.0"
    p.font.size = Pt(88)
    p.font.bold = True
    p.font.color.rgb = NEON_CYAN

    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(2.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True

    p = subtitle_frame.paragraphs[0]
    p.text = "From Prototype to Enterprise Fortification"
    p.font.size = Pt(48)
    p.font.color.rgb = NEON_AMBER
    p.space_before = Pt(12)

    # Add description
    desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(9), Inches(2))
    desc_frame = desc_box.text_frame
    desc_frame.word_wrap = True

    p = desc_frame.paragraphs[0]
    p.text = "Solving Operational Fragility, Latency, and Key Exposure in Autonomous AI Security Gateways"
    p.font.size = Pt(28)
    p.font.color.rgb = TEXT_SECONDARY
    p.line_spacing = 1.3

def add_challenge_solution_slide(prs, title, problem, solution, image_func):
    """Add a challenge/solution slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, DARK_BG)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = NEON_CYAN

    # Problem section
    problem_label = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(4.5), Inches(0.4))
    p = problem_label.text_frame.paragraphs[0]
    p.text = "CHALLENGE"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = NEON_AMBER

    problem_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.85), Inches(4.5), Inches(2.5))
    problem_frame = problem_box.text_frame
    problem_frame.word_wrap = True
    p = problem_frame.paragraphs[0]
    p.text = problem
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_PRIMARY
    p.line_spacing = 1.4

    # Solution section
    solution_label = slide.shapes.add_textbox(Inches(5.2), Inches(1.4), Inches(4.3), Inches(0.4))
    p = solution_label.text_frame.paragraphs[0]
    p.text = "SOLUTION"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = NEON_CYAN

    solution_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.85), Inches(4.3), Inches(2.5))
    solution_frame = solution_box.text_frame
    solution_frame.word_wrap = True
    p = solution_frame.paragraphs[0]
    p.text = solution
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_PRIMARY
    p.line_spacing = 1.4

    # Diagram
    diagram = image_func()
    slide.shapes.add_picture(diagram, Inches(0.5), Inches(4.8), width=Inches(9))

def add_challenge2_slide(prs):
    """Slide 3: High-Availability State Management"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, DARK_BG)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "Challenge 2: High-Availability State Management"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = NEON_CYAN

    # Problem
    problem_label = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(4.5), Inches(0.4))
    p = problem_label.text_frame.paragraphs[0]
    p.text = "CHALLENGE"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = NEON_AMBER

    problem_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.85), Inches(4.5), Inches(2.5))
    problem_frame = problem_box.text_frame
    problem_frame.word_wrap = True
    p = problem_frame.paragraphs[0]
    p.text = "Redis cluster downtime creates a single point of failure that halts the entire security gateway."
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_PRIMARY
    p.line_spacing = 1.4

    # Solution
    solution_label = slide.shapes.add_textbox(Inches(5.2), Inches(1.4), Inches(4.3), Inches(0.4))
    p = solution_label.text_frame.paragraphs[0]
    p.text = "SOLUTION"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = NEON_CYAN

    solution_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.85), Inches(4.3), Inches(2.5))
    solution_frame = solution_box.text_frame
    solution_frame.word_wrap = True
    p = solution_frame.paragraphs[0]
    p.text = "Two-Tier Distributed Nonce Caching: L1 Local Memory LRU Cache + L2 Multi-AZ Redis Cluster Failover"
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_PRIMARY
    p.line_spacing = 1.4

    # Diagram
    diagram = create_cache_flowchart()
    slide.shapes.add_picture(diagram, Inches(0.5), Inches(4.8), width=Inches(9))

def add_challenge3_slide(prs):
    """Slide 4: HITL Bottlenecks"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, DARK_BG)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "Challenge 3: HITL Bottlenecks"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = NEON_CYAN

    # Problem
    problem_label = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(4.5), Inches(0.4))
    p = problem_label.text_frame.paragraphs[0]
    p.text = "CHALLENGE"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = NEON_AMBER

    problem_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.85), Inches(4.5), Inches(2.5))
    problem_frame = problem_box.text_frame
    problem_frame.word_wrap = True
    p = problem_frame.paragraphs[0]
    p.text = "Pausing execution for up to 300s waiting for human approval halts autonomous AI pipelines."
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_PRIMARY
    p.line_spacing = 1.4

    # Solution
    solution_label = slide.shapes.add_textbox(Inches(5.2), Inches(1.4), Inches(4.3), Inches(0.4))
    p = solution_label.text_frame.paragraphs[0]
    p.text = "SOLUTION"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = NEON_CYAN

    solution_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.85), Inches(4.3), Inches(2.5))
    solution_frame = solution_box.text_frame
    solution_frame.word_wrap = True
    p = solution_frame.paragraphs[0]
    p.text = "Async Non-Blocking Webhook Callbacks & Provisional Dry-Runs. Agents continue background tasks while waiting for approval."
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_PRIMARY
    p.line_spacing = 1.4

    # Diagram
    diagram = create_timeline_diagram()
    slide.shapes.add_picture(diagram, Inches(0.5), Inches(4.8), width=Inches(9))

def add_challenge4_slide(prs):
    """Slide 5: Low & Slow Attacks"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, DARK_BG)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "Challenge 4: Multi-Step 'Low & Slow' Attacks"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = NEON_CYAN

    # Problem
    problem_label = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(4.5), Inches(0.4))
    p = problem_label.text_frame.paragraphs[0]
    p.text = "CHALLENGE"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = NEON_AMBER

    problem_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.85), Inches(4.5), Inches(2.5))
    problem_frame = problem_box.text_frame
    problem_frame.word_wrap = True
    p = problem_frame.paragraphs[0]
    p.text = "Stateless AST parsers evaluate commands individually, missing gradual micro-exfiltrations (e.g., 100 x $10 transfers)."
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_PRIMARY
    p.line_spacing = 1.4

    # Solution
    solution_label = slide.shapes.add_textbox(Inches(5.2), Inches(1.4), Inches(4.3), Inches(0.4))
    p = solution_label.text_frame.paragraphs[0]
    p.text = "SOLUTION"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = NEON_CYAN

    solution_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.85), Inches(4.3), Inches(2.5))
    solution_frame = solution_box.text_frame
    solution_frame.word_wrap = True
    p = solution_frame.paragraphs[0]
    p.text = "Stateful Velocity Engine & Sliding-Window Memory. Tracks agent request velocity and cumulative financial thresholds."
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_PRIMARY
    p.line_spacing = 1.4

    # Diagram
    diagram = create_velocity_graph()
    slide.shapes.add_picture(diagram, Inches(0.5), Inches(4.8), width=Inches(9))

def add_challenge5_slide(prs):
    """Slide 6: Private Key Security"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, DARK_BG)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "Challenge 5: Private Key Security"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = NEON_CYAN

    # Problem
    problem_label = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(4.5), Inches(0.4))
    p = problem_label.text_frame.paragraphs[0]
    p.text = "CHALLENGE"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = NEON_AMBER

    problem_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.85), Inches(4.5), Inches(2.5))
    problem_frame = problem_box.text_frame
    problem_frame.word_wrap = True
    p = problem_frame.paragraphs[0]
    p.text = "Storing Ed25519 keys in process memory exposes them to Remote Code Execution (RCE) memory dumps."
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_PRIMARY
    p.line_spacing = 1.4

    # Solution
    solution_label = slide.shapes.add_textbox(Inches(5.2), Inches(1.4), Inches(4.3), Inches(0.4))
    p = solution_label.text_frame.paragraphs[0]
    p.text = "SOLUTION"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = NEON_CYAN

    solution_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.85), Inches(4.3), Inches(2.5))
    solution_frame = solution_box.text_frame
    solution_frame.word_wrap = True
    p = solution_frame.paragraphs[0]
    p.text = "Hardware Security Module (HSM) Integration. AWS KMS / GCP KMS / HashiCorp Vault. Signing keys never touch server memory."
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_PRIMARY
    p.line_spacing = 1.4

    # Diagram
    diagram = create_kms_diagram()
    slide.shapes.add_picture(diagram, Inches(0.5), Inches(4.8), width=Inches(9))

def add_roi_slide(prs):
    """Slide 7: Enterprise ROI & Market Valuation"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, DARK_BG)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "Enterprise ROI & Market Valuation"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = NEON_CYAN

    # Impact Summary Header
    impact_header = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9), Inches(0.4))
    p = impact_header.text_frame.paragraphs[0]
    p.text = "IMPACT SUMMARY"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = NEON_AMBER

    # Impact Content
    impact_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.9), Inches(9), Inches(2.2))
    impact_frame = impact_box.text_frame
    impact_frame.word_wrap = True

    p = impact_frame.paragraphs[0]
    p.text = "• 99.99% uptime SLA achievement"
    p.font.size = Pt(22)
    p.font.color.rgb = TEXT_PRIMARY
    p.space_after = Pt(10)

    p = impact_frame.add_paragraph()
    p.text = "• SOC 2 / ISO 27001 compliance"
    p.font.size = Pt(22)
    p.font.color.rgb = TEXT_PRIMARY
    p.space_after = Pt(10)

    p = impact_frame.add_paragraph()
    p.text = "• $50k–$120k/year enterprise VPC contracts unlocked"
    p.font.size = Pt(22)
    p.font.color.rgb = TEXT_PRIMARY

    # Key Takeaway
    takeaway_header = slide.shapes.add_textbox(Inches(0.5), Inches(4.3), Inches(9), Inches(0.4))
    p = takeaway_header.text_frame.paragraphs[0]
    p.text = "KEY TAKEAWAY"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = NEON_CYAN

    takeaway_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(9), Inches(1.5))
    takeaway_frame = takeaway_box.text_frame
    takeaway_frame.word_wrap = True

    p = takeaway_frame.paragraphs[0]
    p.text = "Resolving operational friction increases compute costs slightly, but elevates market value by "
    p.font.size = Pt(24)
    p.font.color.rgb = TEXT_PRIMARY

    # Add "10x" in amber
    run = p.add_run()
    run.text = "10x"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = NEON_AMBER

    # Comparison Table
    table_header = slide.shapes.add_textbox(Inches(0.5), Inches(6.4), Inches(9), Inches(0.3))
    p = table_header.text_frame.paragraphs[0]
    p.text = "v3 Prototype vs v4 Enterprise Standard"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = NEON_AMBER

    comparison_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(1))
    comparison_frame = comparison_box.text_frame
    comparison_frame.word_wrap = True

    p = comparison_frame.paragraphs[0]
    p.text = "v3: Basic prototype, 95% uptime, manual key management"
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_SECONDARY
    p.space_after = Pt(8)

    p = comparison_frame.add_paragraph()
    p.text = "v4: Enterprise-grade, 99.99% uptime, HSM integration, distributed caching, compliance-ready"
    p.font.size = Pt(18)
    p.font.color.rgb = NEON_CYAN

def add_conclusion_slide(prs):
    """Slide 8: Conclusion & Call to Action"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, DARK_BG)

    # Headline
    headline_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1.2))
    headline_frame = headline_box.text_frame
    headline_frame.word_wrap = True

    p = headline_frame.paragraphs[0]
    p.text = "Unlocking Safe, Autonomous AI Deployment"
    p.font.size = Pt(66)
    p.font.bold = True
    p.font.color.rgb = NEON_CYAN
    p.alignment = PP_ALIGN.CENTER

    # Description
    desc_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1.5))
    desc_frame = desc_box.text_frame
    desc_frame.word_wrap = True

    p = desc_frame.paragraphs[0]
    p.text = "Deploys as an air-gapped Docker container directly inside private enterprise VPCs"
    p.font.size = Pt(32)
    p.font.color.rgb = TEXT_PRIMARY
    p.alignment = PP_ALIGN.CENTER
    p.line_spacing = 1.4

    # Call to Action
    cta_box = slide.shapes.add_textbox(Inches(1.5), Inches(5), Inches(7), Inches(1.5))
    cta_frame = cta_box.text_frame
    cta_frame.word_wrap = True

    p = cta_frame.paragraphs[0]
    p.text = "Schedule an Enterprise Proof-of-Concept (PoC) Today"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = NEON_AMBER
    p.alignment = PP_ALIGN.CENTER

def create_presentation():
    """Create the full presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(prs)

    # Slide 2: Challenge 1 - Token Expiration
    add_challenge_solution_slide(
        prs,
        "Challenge 1: Network Drift & Token Expiration",
        "The 1000ms ephemeral token lifetime causes execution drops during transient network lag or multi-cloud clock drift.",
        "Adaptive 3000ms TTL paired with Atomic JTI Nonce Revocation in Redis. 100% replay protection, zero clock-skew failures.",
        create_comparison_diagram
    )

    # Slide 3: Challenge 2 - HA State Management
    add_challenge2_slide(prs)

    # Slide 4: Challenge 3 - HITL
    add_challenge3_slide(prs)

    # Slide 5: Challenge 4 - Low & Slow
    add_challenge4_slide(prs)

    # Slide 6: Challenge 5 - Key Security
    add_challenge5_slide(prs)

    # Slide 7: ROI & Valuation
    add_roi_slide(prs)

    # Slide 8: Conclusion
    add_conclusion_slide(prs)

    # Save presentation
    output_path = "f:\\Desktop_Data_2026\\Desktop\\INTERLOCK V3\\IntentLock_v4.0_Pitch_Deck.pptx"
    prs.save(output_path)
    print(f"✓ Presentation created: {output_path}")
    return output_path

if __name__ == "__main__":
    create_presentation()
